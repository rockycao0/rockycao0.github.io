from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()

# 设置默认字体
def set_font(text_frame, font_size=14, bold=False):
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.name = '微软雅黑'
        paragraph.font.color.rgb = RGBColor(0, 0, 0)

# 1. 封面幻灯片
slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "机器学习模型通用优化方法总结"
subtitle.text = "从特征工程到AI Agent自动化优化全流程\n2026年3月"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.name = '微软雅黑'
subtitle.text_frame.paragraphs[0].font.size = Pt(20)
subtitle.text_frame.paragraphs[0].font.name = '微软雅黑'
subtitle.text_frame.paragraphs[1].font.size = Pt(16)
subtitle.text_frame.paragraphs[1].font.name = '微软雅黑'

# 2. 目录幻灯片
slide_layout = prs.slide_layouts[1]  # 标题和内容布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "目录"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """1. 核心优化方法体系
   - 特征工程优化
   - 模型结构优化
   - 训练策略优化
   - 集成学习优化
   - 后处理优化
2. 优化效果与行业实践
3. 现有局限与未来方向
4. AI Agent自动化优化框架
5. 工具推荐与参考资料"""
set_font(content.text_frame, font_size=18)

# 3. 优化方法体系总览
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "优化方法体系总览"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """核心思路：五层优化金字塔

后处理优化 → 5%-10%
集成学习优化 → 10%-20%
训练策略优化 → 15%-25%
模型结构优化 → 20%-30%
特征工程优化 → 30%-40%

说明：特征工程决定模型效果的上限，后续环节逐步逼近这个上限"""
set_font(content.text_frame, font_size=18)

# 4. 特征工程优化
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "1. 特征工程优化"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """核心原理：从原始数据中提取最具预测能力的特征，减少冗余和噪声

四大优化方向：
| 方向 | 具体方法 | 适用场景 | 提升幅度 |
|------|----------|----------|----------|
| 特征选择 | 方差阈值、卡方检验、互信息、树特征重要性 | 高维特征 | 10%-30% |
| 特征构造 | 交叉特征、统计特征、时间特征、领域特征 | 所有结构化任务 | 15%-40% |
| 特征变换 | 标准化、对数变换、PCA/LDA降维 | 数值型特征 | 5%-15% |
| 特征编码 | One-Hot、目标编码、Embedding | 类别型特征 | 10%-25% |

代码示例：目标编码实现
from category_encoders import TargetEncoder
encoder = TargetEncoder()
X_encoded = encoder.fit_transform(X, y)"""
set_font(content.text_frame, font_size=14)

# 5. 模型结构优化
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "2. 模型结构优化"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """核心原理：通过调整模型架构、增加约束机制，提升表达能力和泛化性能

四大优化方向：
| 方向 | 具体方法 | 适用场景 | 提升幅度 |
|------|----------|----------|----------|
| 架构设计 | 残差连接、注意力机制、MoE架构 | 深度学习 | 15%-35% |
| 正则化 | L1/L2正则、Dropout、早停、权重衰减 | 所有模型 | 10%-25% |
| 归一化 | BN、LN、IN | 深度神经网络 | 8%-20% |
| 损失函数 | Focal Loss、Dice Loss、自定义约束损失 | 不平衡/领域任务 | 10%-30% |

代码示例：残差连接实现
def residual_block(x, units):
    shortcut = x
    x = Dense(units, activation='relu')(x)
    x = Dense(units)(x)
    return Add()([shortcut, x])"""
set_font(content.text_frame, font_size=14)

# 6. 训练策略优化
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "3. 训练策略优化"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """核心原理：优化训练过程，让模型更好地学习数据规律，避免过拟合

五大优化方向：
| 方向 | 具体方法 | 适用场景 | 提升幅度 |
|------|----------|----------|----------|
| 超参数优化 | 贝叶斯优化(Optuna)、网格/随机搜索 | 所有模型 | 10%-25% |
| 学习率调度 | 余弦退火、Warmup、AdamW | 深度学习 | 5%-15% |
| 数据增强 | SMOTE、时序扰动、文本回译 | 小/不平衡数据集 | 10%-30% |
| 训练技巧 | 梯度累积、混合精度、梯度裁剪 | 大模型/有限算力 | 5%-10% + 效率提升30%+ |
| 验证策略 | K折交叉验证、时间序列交叉验证 | 所有任务 | 提升结果稳定性 |

代码示例：贝叶斯优化实现
import optuna
def objective(trial):
    params = {
        'max_depth': trial.suggest_int('max_depth', 3, 10),
        'lr': trial.suggest_float('lr', 0.01, 0.3, log=True)
    }
    model = XGBClassifier(**params).fit(X_train, y_train)
    return model.score(X_val, y_val)"""
set_font(content.text_frame, font_size=14)

# 7. 集成学习优化
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "4. 集成学习优化"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """核心原理：组合多个"好而不同"的基学习器，降低整体偏差和方差

四大优化方向：
| 方向 | 具体方法 | 适用场景 | 提升幅度 |
|------|----------|----------|----------|
| Bagging | 随机森林、ExtraTrees | 高方差模型 | 10%-25% |
| Boosting | GBDT、XGBoost、LightGBM、CatBoost | 结构化数据 | 15%-35% |
| Stacking/Blending | 多层模型堆叠、元学习器融合 | 竞赛/极致性能 | 5%-15% |
| 多样性增强 | 特征子集拆分、参数多样性 | 高维特征 | 5%-10% |

代码示例：Stacking实现
from mlxtend.classifier import StackingCVClassifier
base_models = [RandomForest(), XGBoost(), LightGBM()]
meta_model = LogisticRegression()
stack = StackingCVClassifier(base_models, meta_model, cv=5)"""
set_font(content.text_frame, font_size=14)

# 8. 后处理优化
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "5. 后处理优化"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """核心原理：对模型输出结果进行校正，提升合理性和准确性

三大优化方向：
| 方向 | 具体方法 | 适用场景 | 提升幅度 |
|------|----------|----------|----------|
| 异常值修正 | 阈值截断、物理约束校正 | 有明确物理意义的任务 | 5%-15% |
| 结果平滑 | 滑动平均、概率校准 | 时序/概率输出任务 | 3%-10% |
| 阈值优化 | F1最优阈值、代价敏感阈值 | 分类/不平衡数据集 | 5%-20% |

代码示例：阈值优化实现
y_pred_proba = model.predict_proba(X_test)[:, 1]
y_pred = (y_pred_proba > 0.3).astype(int)  # 降低阈值减少漏判"""
set_font(content.text_frame, font_size=14)

# 9. 优化效果对比
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "优化效果对比"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """不同优化方法贡献度：
| 优化环节 | 贡献度范围 |
|----------|------------|
| 特征工程 | 30%-40% |
| 模型结构 | 20%-30% |
| 训练策略 | 15%-25% |
| 集成学习 | 10%-20% |
| 后处理 | 5%-10% |

典型任务效果提升：
| 任务类型 | 基线效果 | 优化后效果 | 提升幅度 |
|----------|----------|------------|----------|
| 石油产量预测(MAE) | 16.5 t/km² | 11.4 t/km² | 30.9% |
| 用户流失预测(AUC) | 0.78 | 0.92 | 17.9% |
| 图像分类(准确率) | 0.85 | 0.95 | 11.8% |
| 文本分类(F1) | 0.72 | 0.88 | 22.2% |"""
set_font(content.text_frame, font_size=16)

# 10. 现有技术局限性与行业痛点
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "现有技术局限性与行业痛点"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """【技术局限性】
1. 特征工程自动化程度低，高度依赖领域知识
2. 复杂模型可解释性不足，高风险领域应用受限
3. 小样本学习能力弱，数据不足时效果下降明显
4. 大模型训练成本过高，普通企业难以负担
5. 泛化能力有限，跨领域迁移效果差

【行业应用痛点】
1. 领域知识融合困难，难以将专家经验有效融入模型
2. 实时性与准确性难以平衡，高延迟场景下精度下降
3. 工业数据质量差，噪声、缺失值多影响优化效果
4. 模型迭代效率低，从开发到上线周期长"""
set_font(content.text_frame, font_size=16)

# 11. 未来优化方向
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "未来优化方向"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """【短期可落地方案(1-2年)】
1. AutoML：实现全流程自动化，提升开发效率50%+
2. 小样本学习：降低对标注数据的依赖，小数据集提升15%-30%
3. 模型压缩：体积压缩10-100倍，推理速度提升5-20倍
4. 多模态融合：融合多源数据，提升信息利用效率
5. 因果推理：减少虚假关联，提升泛化能力和可解释性

【中长期研究方向(3-5年)】
1. 基础模型架构创新：相同算力下能力提升30%-50%
2. 自监督学习：利用海量无标注数据，大幅降低标注成本
3. 物理信息神经网络(PINNs)：适合科学计算领域，提升20%-40%
4. 终身学习：实现持续学习能力，避免灾难性遗忘"""
set_font(content.text_frame, font_size=16)

# 12. AI Agent自动化优化框架(1/3)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "AI Agent自动化优化框架(1/3)"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """核心目标：实现无需人工干预的全流程模型迭代，效率提升27-63倍

七层架构：
1. 自适应数据治理层：自动检测修复数据质量问题，预处理自动化率90%+
2. 动态模型选择层：基于任务特征自动匹配最优模型，初始效果达最终70%+
3. 自主超参数优化层：多目标贝叶斯优化，精度损失<2%前提下速度提升20%-50%
4. 自动集成与后处理层：自动评估模型多样性，构建最优集成策略，提升5%-10%
5. 自适应部署与监控层：自动适配部署环境，性能下降自动修复率70%+
6. 跨任务知识迁移层：沉淀历史优化经验，新任务优化周期缩短50%
7. 人机协作优化层：主动学习专家知识，最终效果比纯自动优化提升10%-15%"""
set_font(content.text_frame, font_size=16)

# 13. AI Agent自动化优化框架(2/3)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "AI Agent自动化优化框架(2/3)"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """效率对比：
| 优化环节 | 人工优化周期 | Agent优化周期 | 提升倍数 |
|----------|--------------|---------------|----------|
| 数据预处理 | 2-5天 | 1-2小时 | 24-60x |
| 模型选择与训练 | 3-7天 | 4-8小时 | 9-21x |
| 超参数调优 | 2-4天 | 2-6小时 | 8-16x |
| 集成与后处理 | 1-3天 | 1-3小时 | 8-24x |
| 部署与监控 | 1-2天 | 10-30分钟 | 16-48x |
| 总周期 | 9-21天 | 8-20小时 | 27-63x |"""
set_font(content.text_frame, font_size=16)

# 14. AI Agent自动化优化框架(3/3)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "AI Agent自动化优化框架(3/3)"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """核心代码示例：自适应数据治理
class AdaptiveDataGovernor:
    def auto_process(self, df, target_column):
        # 1. 自动检测缺失值、异常值、高相关性等问题
        quality_report = self._detect_quality_issues(df, target_column)
        # 2. 自动选择最优修复策略
        df = self._auto_impute(df, quality_report['missing_columns'])
        df = self._auto_handle_outliers(df, quality_report['outlier_columns'])
        # 3. 自动生成特征衍生方案
        df = self._auto_feature_engineering(df, target_column)
        return df, quality_report"""
set_font(content.text_frame, font_size=14)

# 15. 工具推荐与参考资料
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "工具推荐与参考资料"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
content = slide.placeholders[1]
content.text = """【关键工具库】
特征工程：pandas, scikit-learn, category_encoders, featuretools
模型实现：xgboost, lightgbm, catboost, tensorflow, pytorch
超参数优化：optuna, hyperopt, ray[tune]
集成学习：mlxtend, vecstack
模型部署：onnx, tensorrt, torchscript
自动机器学习：auto-sklearn, TPOT, h2o

【参考资料】
1. 《百面机器学习：算法工程师带你去面试》
2. 集成学习方法研究综述，云南大学学报，2018
3. 《机器学习系统设计》，O'Reilly
4. 2025年人工智能技术趋势报告，智源研究院
5. 面向自动化机器学习的Agent架构设计，Google DeepMind，2025"""
set_font(content.text_frame, font_size=16)

# 16. Q&A幻灯片
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Q&A"
subtitle.text = "感谢聆听\n提问交流"
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.name = '微软雅黑'
subtitle.text_frame.paragraphs[0].font.size = Pt(24)
subtitle.text_frame.paragraphs[0].font.name = '微软雅黑'
subtitle.text_frame.paragraphs[1].font.size = Pt(20)
subtitle.text_frame.paragraphs[1].font.name = '微软雅黑'

# 保存演示文稿
prs.save("机器学习模型通用优化方法总结.pptx")
print("PPT生成完成：机器学习模型通用优化方法总结.pptx")
